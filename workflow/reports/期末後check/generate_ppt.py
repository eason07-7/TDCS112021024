# -*- coding: utf-8 -*-
"""期末後 check 補充 PPT — 13 頁、聚焦「我們系統怎麼做出來的」

設計原則（per User 要求）：
- 簡約：黑字白底 + 深藍 accent、無動畫、無多餘裝飾
- 封面只有標題、無副標、無組員、無日期
- 13 頁、跟 PDF 章節 1:1 對應（除 3 頁補充頁）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).parent
ASSETS = HERE / "assets"
OUT = HERE / "期末後check_補充簡報.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xD0, 0xD0, 0xD0)
ACCENT = RGBColor(0x2C, 0x5F, 0xA9)
ACCENT_LIGHT = RGBColor(0xE8, 0xEF, 0xF8)
RED = RGBColor(0xC0, 0x39, 0x2B)
FONT = "Microsoft JhengHei"
MONO = "Consolas"


def new_pres():
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def add_text(slide, text, x, y, w, h, *, size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, font=FONT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, lines, x, y, w, h, *, size=14, color=BLACK,
                   font=FONT, line_gap=6, bold_first=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if bold_first and i == 0:
            run.font.bold = True
        p.space_after = Pt(line_gap)
    return tb


def add_rule(slide, y, color=LIGHT_GRAY, height_px=2, x_pad=0.6, width=12.1):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x_pad), y, Inches(width), Emu(height_px * 9525),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_image(slide, path, x, y, w=None, h=None):
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if w is not None:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    if h is not None:
        return slide.shapes.add_picture(str(path), x, y, height=h)
    return slide.shapes.add_picture(str(path), x, y)


def blank_slide(p):
    return p.slides.add_slide(p.slide_layouts[6])


def page_header(slide, page_num, title, subtitle=None):
    """每頁標題列：頁碼 + 標題（左上）+ PDF 對應章節（右上）"""
    add_text(slide, title,
             Inches(0.6), Inches(0.35), Inches(9.5), Inches(0.6),
             size=24, bold=True, color=BLACK)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(11.0), Inches(0.45), Inches(1.8), Inches(0.4),
                 size=11, color=ACCENT, align=PP_ALIGN.RIGHT)
    add_rule(slide, Inches(1.1))


def page_footer(slide, page_num, total=13):
    add_text(slide, "TDCS 國道車流資料下載清洗 CLI",
             Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             size=9, color=GRAY)
    add_text(slide, f"{page_num} / {total}",
             Inches(11.5), Inches(7.05), Inches(1.2), Inches(0.3),
             size=9, color=GRAY, align=PP_ALIGN.RIGHT)


# ── 投影片 1：封面（極簡） ───────────────────────────────────────────
def slide_01_cover(p):
    s = blank_slide(p)
    add_text(s, "TDCS",
             Inches(0.6), Inches(2.4), Inches(12.1), Inches(1.0),
             size=24, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, "國道車流資料下載清洗 CLI",
             Inches(0.6), Inches(3.0), Inches(12.1), Inches(1.8),
             size=48, bold=True, color=BLACK, align=PP_ALIGN.CENTER)


# ── 投影片 2：問題提出（放 PDF §2.1 截圖）─────────────────────────
def slide_02_problem(p):
    s = blank_slide(p)
    page_header(s, 2, "書面報告中的編號為何看不懂？", "對應 PDF §2.1")

    add_text(s, "書面報告 §2.1 配置一覽表右欄 — 讀者看到「demo_17 / demo_18 ...」無法對應到任何東西",
             Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.5),
             size=14, color=BLACK)

    # PDF §2.1 截圖
    add_image(s, ASSETS / "pdf_section_21.png",
              Inches(1.5), Inches(2.0), w=Inches(10.3))

    add_text(s, "「配圖」欄出現 demo_17 / demo_18，但 PDF 完整看不到 demo_17 是什麼",
             Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.4),
             size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    page_footer(s, 2)


# ── 投影片 3：原因 + 6 張 VS Code proof ─────────────────────────────
def slide_03_proof(p):
    s = blank_slide(p)
    page_header(s, 3, "原因：開發代號忘記翻成 Figure 編號", "補充章")

    add_text(s, "這 6 張是 VS Code 開我們專案的截圖、左欄檔名真的就叫 demo_XX_*.png — 內部代號、不是事後虛構",
             Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.4),
             size=12, color=GRAY)

    proofs = [
        ("proof_vscode_01_demo17_lambda.png", "demo_17 → Lambda 配置"),
        ("proof_vscode_02_demo18_sqs.png", "demo_18 → SQS 佇列"),
        ("proof_vscode_03_demo12_athena.png", "demo_12 → Athena workgroup"),
        ("proof_vscode_04_demo19_glue.png", "demo_19 → Glue 資料表"),
        ("proof_vscode_05_demo20_apigw.png", "demo_20 → API Gateway 路由"),
        ("proof_vscode_06_demo21_s3.png", "demo_21 → S3 資料夾結構"),
    ]
    cols, rows = 3, 2
    grid_x = Inches(0.5)
    grid_y = Inches(1.85)
    cell_w = Inches(4.1)
    cell_h = Inches(2.0)
    gap_x = Inches(0.05)
    gap_y = Inches(0.5)

    for i, (fn, label) in enumerate(proofs):
        col = i % cols
        row = i // cols
        x = grid_x + (cell_w + gap_x) * col
        y = grid_y + (cell_h + gap_y) * row
        add_image(s, ASSETS / fn, x, y, w=cell_w, h=cell_h)
        add_text(s, label, x, y + cell_h + Inches(0.05),
                 cell_w, Inches(0.4),
                 size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text(s, "實作組員用代號傳給文書組員、文書組員整合時忘了改成 Figure 編號 — 流程疏失、非造假",
             Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.35),
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    page_footer(s, 3)


# ── 投影片 4：整體架構 ──────────────────────────────────────────────
def slide_04_architecture(p):
    s = blank_slide(p)
    page_header(s, 4, "系統整體架構 — 7 步端到端流程", "對應 PDF §2")

    add_image(s, ASSETS / "系統架構圖.png",
              Inches(0.6), Inches(1.3), w=Inches(12.1))

    add_text(s, "(1) 本機下載 → (2) 上傳 S3 → (3) POST /clean → (4) 進 SQS 排隊 → (5) Lambda 清洗 → (6) Glue 註冊 → (7) Athena 查詢",
             Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.4),
             size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    page_footer(s, 4)


# ── 投影片 5：本地 CLI - 技術選擇 ────────────────────────────────────
def slide_05_cli_stack(p):
    s = blank_slide(p)
    page_header(s, 5, "本地端 CLI：技術選擇", "對應 PDF §2.2")

    # 左半：選了什麼
    add_text(s, "技術棧",
             Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.5),
             size=18, bold=True, color=ACCENT)
    add_paragraphs(s, [
        "• 語言：TypeScript（不用 Python）",
        "• 終端 UI：ink（在終端機跑 React）",
        "• 打包：esbuild → 單一 dist/index.js",
        "• 入口：bin/tdcs-dl、可 npm 全域裝",
        "• AWS SDK：v3（modular、tree-shake）",
    ], Inches(0.6), Inches(2.1), Inches(5.5), Inches(3.5),
        size=14, color=BLACK, line_gap=10)

    # 右半：為什麼
    add_text(s, "為什麼這樣選",
             Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
             size=18, bold=True, color=ACCENT)
    add_paragraphs(s, [
        "• TypeScript：npm 一行就能裝、不必處理 Python venv、跨平台 OK",
        "• ink：互動式選單比 prompt 一行一行問友善 10 倍",
        "• esbuild 單檔：使用者不必裝任何依賴、節省磁碟與環境麻煩",
        "• 整套都是現成成熟工具、不重新發明輪子",
    ], Inches(7.0), Inches(2.1), Inches(5.7), Inches(3.5),
        size=13, color=BLACK, line_gap=10)

    add_text(s, "課堂沒教的工具：ink / esbuild — 自學官方文件 + 範例專案",
             Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    page_footer(s, 5)


# ── 投影片 6：本地 CLI - 五步 wizard 設計 ────────────────────────────
def slide_06_wizard(p):
    s = blank_slide(p)
    page_header(s, 6, "本地端 CLI：五步 wizard 怎麼設計", "對應 PDF §2.2")

    # 5 步格子
    steps = [
        ("Step 1", "資料類型", "選 M06A"),
        ("Step 2", "時間區間", "整月 / 單日"),
        ("Step 3", "交流道門架", "多選"),
        ("Step 4", "輸出位置", "S3 / 本機"),
        ("Step 5", "確認", "Enter 送出"),
    ]
    box_w = Inches(2.3)
    box_h = Inches(1.5)
    gap = Inches(0.2)
    total_w = box_w * 5 + gap * 4
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(1.6)

    for i, (tag, title, desc) in enumerate(steps):
        x = start_x + (box_w + gap) * i
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, box_w, box_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = ACCENT_LIGHT
        rect.line.color.rgb = ACCENT
        tf = rect.text_frame
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(10)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = tag
        run.font.name = FONT; run.font.size = Pt(13); run.font.bold = True
        run.font.color.rgb = ACCENT
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = title
        r2.font.name = FONT; r2.font.size = Pt(17); r2.font.bold = True
        r2.font.color.rgb = BLACK
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run(); r3.text = desc
        r3.font.name = FONT; r3.font.size = Pt(11)
        r3.font.color.rgb = GRAY

    # 設計理由
    add_text(s, "設計理由",
             Inches(0.6), Inches(3.5), Inches(12.1), Inches(0.5),
             size=18, bold=True, color=ACCENT)
    add_paragraphs(s, [
        "• 每步只一個決策 — 不一次壓垮使用者",
        "• ESC 可隨時回上一步 — 容錯設計、不必重來",
        "• 狀態用 state machine 管：goNext / goPrev / state.answers",
        "• 5 個 React 元件（Data.tsx / Time.tsx / Gantry.tsx / Output.tsx / Confirm.tsx）— 模組化、可獨立測試",
        "• 確認頁顯示所有選擇後才實際送出 — 提交前最後檢查",
    ], Inches(0.6), Inches(4.1), Inches(12.1), Inches(2.5),
        size=13, color=BLACK, line_gap=8)

    page_footer(s, 6)


# ── 投影片 7：下載策略 E6 ────────────────────────────────────────────
def slide_07_e6(p):
    s = blank_slide(p)
    page_header(s, 7, "下載策略 E6：7 種候選方案的科學取捨", "對應 PDF §3.6")

    add_text(s, "核心問題：TDCS 公開資料只認台灣 IP、我們的 AWS 卻鎖在美國東部 — 怎麼讓 us-east-1 抓到台灣 IP 限定的資料？",
             Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.5),
             size=13, color=BLACK)

    # 表格
    rows = [
        ["#", "候選方案", "狀態", "結論"],
        ["E1", "Lambda 直連 TDCS", "✅ 實測：timeout 15,990ms", "❌ 確認封鎖"],
        ["E2", "EC2 直連 TDCS", "✅ 實測：curl exit 28", "❌ 同 E1、區域層級封鎖"],
        ["E3", "第三方 proxy（BrightData 等）", "📋 文件分析", "需付費 + 信用卡"],
        ["E4", "Tailscale + EC2 subnet router", "📋 架構分析", "3 SPOF 風險"],
        ["E5", "nginx + Cloudflare tunnel", "📋 設計分析", "免費但 URL 不穩"],
        ["E6", "純 client-side（本機下載 → S3）", "✅ PoC 通過 + 麻豆段 22 GB 真跑", "✅ Winner"],
        ["E7", "Hybrid S3 cache + webhook", "📋 架構設計", "依賴鏈太複雜"],
    ]

    col_widths = [Inches(0.8), Inches(4.6), Inches(4.0), Inches(3.0)]
    row_h = Inches(0.5)
    table_x = Inches(0.45)
    table_y = Inches(2.0)

    for r_idx, row in enumerate(rows):
        x = table_x
        for c_idx, cell in enumerate(row):
            w = col_widths[c_idx]
            y = table_y + row_h * r_idx
            is_header = r_idx == 0
            is_winner = r_idx == 6
            cell_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, row_h)
            if is_header:
                cell_box.fill.solid()
                cell_box.fill.fore_color.rgb = ACCENT
            elif is_winner:
                cell_box.fill.solid()
                cell_box.fill.fore_color.rgb = ACCENT_LIGHT
            else:
                cell_box.fill.solid()
                cell_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell_box.line.color.rgb = LIGHT_GRAY
            tf = cell_box.text_frame
            tf.margin_left = Pt(6)
            tf.margin_right = Pt(6)
            tf.margin_top = Pt(4)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c_idx != 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = cell
            run.font.name = FONT
            run.font.size = Pt(10)
            run.font.bold = is_header or is_winner
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if is_header else BLACK
            x += w

    add_text(s, "結論：下載必須走本機（台灣 IP 是物理約束）、AWS 只負責清洗 + 查詢 — 對的東西在對的地方",
             Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
             size=12, italic=True, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    page_footer(s, 7)


# ── 投影片 8：CLI ↔ AWS 對接 ────────────────────────────────────────
def slide_08_integration(p):
    s = blank_slide(p)
    page_header(s, 8, "CLI ↔ AWS 怎麼對接", "對應 PDF §3.2")

    # 左：endpoint 三層
    add_text(s, "Endpoint 三層 priority",
             Inches(0.6), Inches(1.4), Inches(5.7), Inches(0.5),
             size=16, bold=True, color=ACCENT)
    add_paragraphs(s, [
        "1. env 變數（TDCS_DL_ENDPOINT） — 開發測試用",
        "2. config 檔（~/.tdcs-dl/config.json） — 自架者用",
        "3. 預設 hardcoded — 本 demo 帳號用",
        "",
        "為什麼三層：",
        "→ 教學版預設 hardcode、不用設定",
        "→ 自架者可改 config 指自己的 API GW",
        "→ 同一份 CLI 兼容兩種使用情境",
    ], Inches(0.6), Inches(2.0), Inches(5.7), Inches(4.5),
        size=12, color=BLACK, line_gap=6)

    # 右：POST + polling
    add_text(s, "觸發 + 輪詢機制",
             Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.5),
             size=16, bold=True, color=ACCENT)
    add_paragraphs(s, [
        "1. CLI POST /clean → API GW 秒回 job_id",
        "2. CLI 每 5 秒 GET /jobs/{id} 一次",
        "3. 狀態流：accepted → processing → done",
        "4. 看到 done 才結束等待",
        "",
        "為什麼這樣：",
        "→ POST 立刻回、不卡 30 秒 timeout",
        "→ Lambda 慢慢清、CLI 慢慢問",
        "→ HTTP 標準、不用建長連線",
    ], Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.5),
        size=12, color=BLACK, line_gap=6)

    add_text(s, "三層 endpoint 是「demo 版 vs 自架版」共用一份 code 的關鍵設計",
             Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    page_footer(s, 8)


# ── 投影片 9：Lambda Container ─────────────────────────────────────
def slide_09_lambda(p):
    s = blank_slide(p)
    page_header(s, 9, "Lambda Container：為什麼用 image 不用 zip", "對應 PDF §3.3 + §3.4 + §3.5")

    add_paragraphs(s, [
        "Lambda 部署有兩種方式：",
        "• Zip 包：上傳程式碼壓縮檔、限 250 MB、只能用預設 runtime",
        "• Container image：上傳 Docker image、限 10 GB、可自帶 runtime",
    ], Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.5),
        size=14, color=BLACK, line_gap=8)

    add_text(s, "為什麼我們選 Container",
             Inches(0.6), Inches(3.0), Inches(12.1), Inches(0.5),
             size=16, bold=True, color=ACCENT)

    add_paragraphs(s, [
        "1. 清洗用 nodejs-polars（Rust 寫的 native 模組）— 不在 Lambda 預設 runtime 裡",
        "2. Container 讓我們用 Docker 把 polars + Node.js 包成 image、推 ECR、Lambda 拉來跑",
        "3. Image tag = git commit SHA — 每次部署 tag 都不同、Terraform 看得到變化會更新",
        "4. 配置：memory 2048 MB / timeout 900 s / ephemeral 1024 MB（§2.1 配置表）",
    ], Inches(0.6), Inches(3.6), Inches(12.1), Inches(2.8),
        size=13, color=BLACK, line_gap=10)

    add_text(s, "git SHA tag 解決「:latest 鎖死、Lambda 不更新」的鬼打牆問題（§3.3）",
             Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    page_footer(s, 9)


# ── 投影片 10：SQS broker ──────────────────────────────────────────
def slide_10_sqs(p):
    s = blank_slide(p)
    page_header(s, 10, "SQS broker：30 秒對上 14 分鐘的鴻溝", "對應 PDF §3.2")

    # 問題
    add_text(s, "問題：API Gateway 最多等 30 秒、Lambda 清洗要 847 秒（14 分鐘）",
             Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.5),
             size=14, bold=True, color=RED)

    # 對比表
    rows = [
        ["", "限制", "我們需要"],
        ["API Gateway", "30 秒 timeout", "—"],
        ["Lambda", "15 分鐘 timeout", "清洗實測 847 秒"],
        ["差距", "—", "28 倍"],
    ]
    col_w = [Inches(2.8), Inches(4.5), Inches(4.5)]
    row_h = Inches(0.55)
    tx = Inches(0.75); ty = Inches(2.05)
    for ri, row in enumerate(rows):
        x = tx
        for ci, cell in enumerate(row):
            box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, ty + row_h * ri, col_w[ci], row_h)
            box.fill.solid()
            box.fill.fore_color.rgb = ACCENT if ri == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            box.line.color.rgb = LIGHT_GRAY
            tf = box.text_frame
            tf.margin_left = Pt(8); tf.margin_top = Pt(6)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = cell
            r.font.name = FONT; r.font.size = Pt(12)
            r.font.bold = ri == 0 or ci == 0
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else BLACK
            x += col_w[ci]

    # 解法
    add_text(s, "解法：SQS 排隊 + producer/consumer 拆兩段",
             Inches(0.6), Inches(4.6), Inches(12.1), Inches(0.5),
             size=16, bold=True, color=ACCENT)

    add_paragraphs(s, [
        "• Producer 路徑（API GW 觸發）：收到請求 → 寫 SQS → 立刻回 job_id（< 1 秒）",
        "• Consumer 路徑（SQS event source）：背景做 14 分鐘清洗、不必任何人等",
        "• 同一個 Lambda 兩個 trigger — 程式碼一份、兩個進入點",
        "• DLQ 死信佇列 maxReceiveCount=2：失敗重試 1 次後送死信、不無限重跑",
    ], Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.0),
        size=12, color=BLACK, line_gap=6)

    page_footer(s, 10)


# ── 投影片 11：Athena + Glue ───────────────────────────────────────
def slide_11_athena(p):
    s = blank_slide(p)
    page_header(s, 11, "Athena + Glue：為什麼選 serverless 查詢", "對應 PDF §3.5 + §4.1")

    # 左：Glue
    add_text(s, "Glue Data Catalog",
             Inches(0.6), Inches(1.4), Inches(5.7), Inches(0.5),
             size=16, bold=True, color=ACCENT)
    add_paragraphs(s, [
        "• 註冊 cleaned_v2_skeleton 資料表",
        "• 9 欄 snake_case + partition yyyymm",
        "• Lambda 寫 Parquet 完跑 MSCK REPAIR",
        "• 自動發現新 partition、不必手動加",
    ], Inches(0.6), Inches(2.0), Inches(5.7), Inches(2.5),
        size=13, color=BLACK, line_gap=8)

    # 右：Athena
    add_text(s, "Athena",
             Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.5),
             size=16, bold=True, color=ACCENT)
    add_paragraphs(s, [
        "• Serverless SQL：開查就用、沒在用就 $0",
        "• 按 scan 大小計費（$5 / TB）",
        "• Workgroup 設 10 MB scan cap 護網",
        "• 實測查 1 次掃 50-100 KB、計費 $0.0000005",
    ], Inches(6.8), Inches(2.0), Inches(6.0), Inches(2.5),
        size=13, color=BLACK, line_gap=8)

    # 為什麼不用 RDS / RedShift
    add_text(s, "為什麼不用 RDS / RedShift？",
             Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.5),
             size=16, bold=True, color=ACCENT)
    add_paragraphs(s, [
        "• 那是「常駐機器」、按月計費（RDS 約 $50/月、RedShift 約 $180/月）",
        "• 我們資料量小（17 GB）、查詢稀疏（一天幾次）— 用不到常駐",
        "• Athena 規模小時就是免費、規模大時才開始付費 — 完全 fit 學期專題",
    ], Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.8),
        size=12, color=BLACK, line_gap=8)

    page_footer(s, 11)


# ── 投影片 12：IaC + Terraform ─────────────────────────────────────
def slide_12_iac(p):
    s = blank_slide(p)
    page_header(s, 12, "IaC：所有 AWS 設定都用程式碼管", "對應 PDF §2.1 雙錨對齊")

    add_text(s, "infra/terraform/ 共 11 個 .tf 檔案",
             Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.5),
             size=16, bold=True, color=ACCENT)

    # 11 個 .tf 列表（grid）
    tfs = [
        ("main.tf", "AWS provider、region 鎖定"),
        ("variables.tf", "Account ID、bucket 名等變數"),
        ("s3.tf", "S3 bucket 引用 + 4 個 prefix"),
        ("lambda.tf", "Lambda function tdcs-dl-cleaner"),
        ("apigw.tf", "API Gateway HTTP API + 2 routes"),
        ("sqs.tf", "SQS queue + DLQ + event source mapping"),
        ("glue.tf", "Glue database + table + 9 欄 schema"),
        ("athena.tf", "Athena workgroup + 10 MB scan cap"),
        ("budget.tf", "F-H3 cost guard"),
        ("outputs.tf", "api_gw_url / lambda_arn 等輸出"),
        ("versions.tf", "Terraform 與 provider 版本鎖定"),
    ]
    grid_x = Inches(0.6); grid_y = Inches(2.0)
    cell_w = Inches(6.0); cell_h = Inches(0.4)
    for i, (name, desc) in enumerate(tfs):
        col = i // 6; row = i % 6
        x = grid_x + cell_w * col + Inches(0.1) * col
        y = grid_y + cell_h * row
        add_text(s, name, x, y, Inches(2.0), cell_h, size=12, bold=True,
                 font=MONO, color=ACCENT)
        add_text(s, desc, x + Inches(2.0), y, Inches(4.0), cell_h, size=11,
                 color=BLACK)

    # 意義
    add_text(s, "意義",
             Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.5),
             size=16, bold=True, color=ACCENT)
    add_paragraphs(s, [
        "• 任何人 clone 這個 repo、在自己 AWS 帳號跑 `terraform apply` → 重建一模一樣的後端",
        "• 配置是程式碼、不是「在 Console 手點完留個截圖」 — Console live 值 = .tf 程式碼，零漂移（§2.1 雙錨）",
        "• 這是 v3.0「人人可自架」的基礎 — 不必跟我們拿金鑰、用自己的帳號就行",
    ], Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.8),
        size=12, color=BLACK, line_gap=8)

    page_footer(s, 12)


# ── 投影片 13：Live demo 入口 ────────────────────────────────────────
def slide_13_demo(p):
    s = blank_slide(p)

    add_text(s, "現場重現一次",
             Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.8),
             size=28, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # 4 步流程
    steps = [
        ("1", "AWS 確認", "看 S3 raw 目前只有 3 / 4 / 5 月、沒有 6 月"),
        ("2", "本機跑 CLI", "node cli/dist/index.js → 五步精靈 → 選 0601 單日"),
        ("3", "看 S3 raw 新檔", "yyyymm=202606/ 出現、含 24 個 csv.gz"),
        ("4", "看清洗結果", "cleaned.parquet + Athena 查 — 數字對得上選的範圍"),
    ]

    y = Inches(2.8)
    for tag, title, desc in steps:
        # 圓圈標號
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y, Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        tf = circle.text_frame
        tf.margin_left = 0; tf.margin_top = Pt(2)
        para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
        run = para.add_run(); run.text = tag
        run.font.name = FONT; run.font.size = Pt(20); run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        add_text(s, title,
                 Inches(2.6), y + Inches(0.05), Inches(3.0), Inches(0.5),
                 size=18, bold=True, color=BLACK)
        add_text(s, desc,
                 Inches(5.5), y + Inches(0.1), Inches(7.0), Inches(0.5),
                 size=14, color=GRAY)
        y += Inches(0.85)

    add_text(s, "請看螢幕切到 AWS Console →",
             Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.6),
             size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text(s, "github.com/eason07-7/TDCS112021024  ‧  所有 code、配置、截圖均公開",
             Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3),
             size=10, color=GRAY, align=PP_ALIGN.CENTER, font=MONO)


def main():
    p = new_pres()
    slide_01_cover(p)
    slide_02_problem(p)
    slide_03_proof(p)
    slide_04_architecture(p)
    slide_05_cli_stack(p)
    slide_06_wizard(p)
    slide_07_e6(p)
    slide_08_integration(p)
    slide_09_lambda(p)
    slide_10_sqs(p)
    slide_11_athena(p)
    slide_12_iac(p)
    slide_13_demo(p)
    p.save(str(OUT))
    print(f"saved: {OUT}")
    print(f"size: {OUT.stat().st_size:,} bytes")
    print(f"slides: {len(p.slides)}")


if __name__ == "__main__":
    main()
